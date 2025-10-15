from app.database.connection import SessionLocal
from app.database.crud import FormularioCRUD
import streamlit as st
import pandas as pd
from datetime import datetime, date, timedelta
import sys
import os

# Add the project root to the path
sys.path.insert(0, os.path.dirname(os.path.dirname(
    os.path.dirname(os.path.abspath(__file__)))))


# Try to import optional components with error handling
try:
    from dashboard.components.interactive_filters import InteractiveFilters
    FILTERS_AVAILABLE = True
except Exception as e:
    FILTERS_AVAILABLE = False

try:
    from app.utils.report_history import ReportHistory
    HISTORY_AVAILABLE = True
except Exception as e:
    HISTORY_AVAILABLE = False

try:
    from app.utils.report_generator import ReportGenerator
    GENERATOR_AVAILABLE = True
except Exception as e:
    GENERATOR_AVAILABLE = False


def show_report_generation_page():
    """Report generation page with NLG capabilities"""

    # Require authentication
    from app.auth.streamlit_auth import auth
    if not auth.is_authenticated():
        auth.show_login_form()
        return

    st.title("📄 Generación de Reportes")
    st.markdown("Genere reportes automáticos de actividades docentes.")

    # Add cache refresh button
    col1, col2, col3 = st.columns([1, 1, 4])
    with col1:
        if st.button("🔄 Actualizar Datos", help="Limpiar cache y recargar datos"):
            st.cache_data.clear()
            st.rerun()
    with col2:
        st.write(
            f"*Última actualización: {datetime.now().strftime('%H:%M:%S')}*")

    # Check component availability
    components_missing = []
    if not GENERATOR_AVAILABLE:
        components_missing.append("Generador de reportes")
    if not HISTORY_AVAILABLE:
        components_missing.append("Historial de reportes")
    if not FILTERS_AVAILABLE:
        components_missing.append("Filtros interactivos")

    if components_missing:
        st.warning(
            f"⚠️ Algunos componentes no están disponibles: {', '.join(components_missing)}")
        st.info("💡 Usando funcionalidad básica de reportes.")

    # Initialize available components
    report_generator = None
    history_manager = None
    filters = None

    if GENERATOR_AVAILABLE:
        try:
            report_generator = ReportGenerator()
        except Exception as e:
            st.warning(f"No se pudo inicializar el generador: {e}")

    if HISTORY_AVAILABLE:
        try:
            history_manager = ReportHistory()
        except Exception as e:
            st.warning(f"No se pudo inicializar el historial: {e}")

    if FILTERS_AVAILABLE:
        try:
            filters = InteractiveFilters()
        except Exception as e:
            st.warning(f"No se pudieron inicializar los filtros: {e}")

    # Load data
    try:
        all_forms, metrics = load_report_data()
    except Exception as e:
        st.error(f"Error al cargar datos: {e}")
        return

    if not all_forms:
        st.warning("No hay datos disponibles para generar reportes.")
        return

    # Sidebar configuration
    st.sidebar.header("🎛️ Configuración de Reporte")

    # Report type selection
    report_type = st.sidebar.selectbox(
        "Tipo de reporte:",
        ["annual", "quarterly", "data_table"],
        format_func=lambda x: {
            "annual": "📊 Reporte Anual Narrativo",
            "quarterly": "📈 Resumen Trimestral",
            "data_table": "📋 Reporte de Datos Tabulares"
        }[x]
    )

    # Period selection
    if report_type == "quarterly":
        col1, col2 = st.sidebar.columns(2)
        with col1:
            year = st.selectbox("Año:", range(2020, 2051),
                                index=5)  # Default to 2025
        with col2:
            quarter = st.selectbox("Trimestre:", [1, 2, 3, 4])
    elif report_type == "annual":
        # For annual reports, just select the year
        year = st.sidebar.selectbox("Año del reporte:", range(2020, 2051),
                                    index=5)  # Default to 2025
    else:
        # Date range for custom reports only
        min_date = min(
            [f.fecha_envio for f in all_forms if f.fecha_envio]).date()
        max_date = max(
            [f.fecha_envio for f in all_forms if f.fecha_envio]).date()

        date_range = st.sidebar.date_input(
            "Período del reporte:",
            value=(min_date, max_date),
            min_value=min_date,
            max_value=max_date
        )

    # Additional options
    with st.sidebar.expander("⚙️ Opciones Avanzadas"):
        include_trends = st.checkbox(
            "Incluir análisis de tendencias", value=True)
        include_highlights = st.checkbox(
            "Incluir elementos destacados", value=True)
        report_tone = st.selectbox(
            "Tono del reporte:",
            ["professional", "academic", "executive"],
            format_func=lambda x: {
                "professional": "Profesional",
                "academic": "Académico",
                "executive": "Ejecutivo"
            }[x]
        )

        custom_title = st.text_input(
            "Título personalizado (opcional):",
            placeholder="Reporte de Actividades Docentes"
        )

    # Main content area
    if history_manager:
        tab1, tab2, tab3 = st.tabs(
            ["🔧 Generar Reporte", "📚 Historial", "📊 Estadísticas"])
    else:
        tab1 = st.container()
        st.info("💡 Funciones de historial no disponibles - solo generación básica")

    with tab1:
        # Report generation interface
        st.subheader("Generación de Reporte")

        # Filter data based on period
        if report_type == "quarterly":
            period_start, period_end = get_quarter_dates(quarter, year)
            filtered_forms = filter_forms_by_period(
                all_forms, period_start, period_end)
            period_text = f"{get_quarter_name(quarter)} {year}"
        elif report_type == "annual":
            # For annual reports, use the entire year
            period_start = date(year, 1, 1)
            period_end = date(year, 12, 31)
            filtered_forms = filter_forms_by_period(
                all_forms, period_start, period_end)
            period_text = f"Año {year}"
        else:
            if len(date_range) == 2:
                period_start, period_end = date_range
                filtered_forms = filter_forms_by_period(
                    all_forms, period_start, period_end)
                period_text = f"{period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}"
            else:
                st.warning("Por favor seleccione un rango de fechas válido.")
                return

        # Show preview information
        col1, col2, col3 = st.columns(3)

        with col1:
            st.metric("Período", period_text)
        with col2:
            st.metric("Formularios", len(filtered_forms))
        with col3:
            approved_count = len(
                [f for f in filtered_forms if f.estado.value == 'APROBADO'])
            st.metric("Aprobados", approved_count)

        if filtered_forms:
            # Show data preview
            with st.expander("👀 Vista Previa de Datos"):
                preview_df = create_preview_dataframe(filtered_forms)

                # Add pagination controls
                total_records = len(preview_df)
                st.write(f"**Total de registros:** {total_records}")

                if total_records > 0:
                    # Pagination
                    records_per_page = st.selectbox("Registros por página:", [
                                                    10, 25, 50, 100], index=0)
                    total_pages = (total_records - 1) // records_per_page + 1

                    if total_pages > 1:
                        page = st.selectbox(
                            "Página:", range(1, total_pages + 1))
                        start_idx = (page - 1) * records_per_page
                        end_idx = start_idx + records_per_page
                        display_df = preview_df.iloc[start_idx:end_idx]
                        st.write(
                            f"Mostrando registros {start_idx + 1} - {min(end_idx, total_records)} de {total_records}")
                    else:
                        display_df = preview_df

                    st.dataframe(display_df, use_container_width=True)
                else:
                    st.info(
                        "No hay registros para mostrar en el período seleccionado.")

                # Activity summary
                st.subheader("Resumen de Actividades")
                activity_summary = calculate_activity_summary(filtered_forms)

                col1, col2, col3 = st.columns(3)
                with col1:
                    st.metric("Cursos", activity_summary['cursos'])
                    st.metric("Publicaciones",
                              activity_summary['publicaciones'])
                with col2:
                    st.metric("Eventos", activity_summary['eventos'])
                    st.metric("Diseños Curriculares",
                              activity_summary['disenos'])
                with col3:
                    st.metric("Movilidades", activity_summary['movilidades'])
                    st.metric("Reconocimientos",
                              activity_summary['reconocimientos'])

                # Segunda fila para certificaciones
                col4, col5, col6 = st.columns(3)
                with col4:
                    st.metric("Certificaciones",
                              activity_summary['certificaciones'])

            # Export format selection
            st.subheader("📤 Opciones de Exportación")

            col1, col2 = st.columns(2)

            with col1:
                # Generate markdown report
                if st.button("🚀 Generar Reporte Markdown", type="primary"):
                    if report_generator and history_manager:
                        generate_and_display_report(
                            report_generator, history_manager, filtered_forms,
                            report_type, period_start, period_end,
                            include_trends, include_highlights, report_tone, custom_title
                        )
                    else:
                        # Generate simple report without advanced components
                        generate_simple_report_display(
                            filtered_forms, report_type, period_start, period_end,
                            custom_title or f"Reporte {report_type} {period_start.year}"
                        )

            with col2:
                # Multi-format export options
                export_format = st.selectbox(
                    "Formato de exportación:",
                    ["PDF", "Excel", "PowerPoint"],
                    help="Seleccione el formato para exportar el reporte"
                )

                # Generate content for direct download
                title = custom_title or f"Reporte {report_type} {period_start.year}"
                filename_base = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}"

                if export_format == "PDF":
                    try:
                        from reportlab.lib.pagesizes import A4
                        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                        from reportlab.lib import colors
                        from io import BytesIO

                        buffer = BytesIO()
                        doc = SimpleDocTemplate(buffer, pagesize=A4)
                        styles = getSampleStyleSheet()
                        story = []

                        # Generate report content
                        report_content = generate_simple_report(
                            filtered_forms, title, report_type, period_start, period_end)

                        # Convert to PDF format
                        lines = report_content.split('\n')
                        for line in lines:
                            if line.strip():
                                if line.startswith('# '):
                                    story.append(
                                        Paragraph(line[2:], styles['Heading1']))
                                elif line.startswith('## '):
                                    story.append(
                                        Paragraph(line[3:], styles['Heading2']))
                                else:
                                    story.append(
                                        Paragraph(line, styles['Normal']))
                                story.append(Spacer(1, 6))

                        doc.build(story)
                        buffer.seek(0)
                        pdf_content = buffer.getvalue()

                        st.download_button(
                            label=f"📄 Exportar como {export_format}",
                            data=pdf_content,
                            file_name=f"{filename_base}.pdf",
                            mime="application/pdf",
                            key=f"direct_pdf_{datetime.now().timestamp()}"
                        )
                    except ImportError:
                        report_content = generate_simple_report(
                            filtered_forms, title, report_type, period_start, period_end)
                        st.download_button(
                            label=f"📄 Exportar como {export_format}",
                            data=report_content,
                            file_name=f"{filename_base}.txt",
                            mime="text/plain",
                            key=f"direct_txt_{datetime.now().timestamp()}"
                        )

                elif export_format == "Excel":
                    try:
                        from openpyxl import Workbook
                        from openpyxl.styles import Font
                        from io import BytesIO

                        wb = Workbook()
                        ws = wb.active
                        ws.title = "Reporte"

                        # Generate report content
                        report_content = generate_simple_report(
                            filtered_forms, title, report_type, period_start, period_end)

                        # Add content to Excel
                        row = 1
                        lines = report_content.split('\n')
                        for line in lines:
                            if line.strip():
                                ws.cell(row=row, column=1, value=line)
                                if line.startswith('# '):
                                    ws.cell(row=row, column=1).font = Font(
                                        size=18, bold=True)
                                elif line.startswith('## '):
                                    ws.cell(row=row, column=1).font = Font(
                                        size=14, bold=True)
                                row += 1

                        ws.column_dimensions['A'].width = 80

                        buffer = BytesIO()
                        wb.save(buffer)
                        buffer.seek(0)
                        excel_content = buffer.getvalue()

                        st.download_button(
                            label=f"📄 Exportar como {export_format}",
                            data=excel_content,
                            file_name=f"{filename_base}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            key=f"direct_excel_{datetime.now().timestamp()}"
                        )
                    except ImportError:
                        report_content = generate_simple_report(
                            filtered_forms, title, report_type, period_start, period_end)
                        st.download_button(
                            label=f"📄 Exportar como {export_format}",
                            data=report_content,
                            file_name=f"{filename_base}.csv",
                            mime="text/csv",
                            key=f"direct_csv_{datetime.now().timestamp()}"
                        )

                elif export_format == "PowerPoint":
                    try:
                        from pptx import Presentation
                        from io import BytesIO

                        prs = Presentation()

                        # Title slide
                        slide_layout = prs.slide_layouts[0]
                        slide = prs.slides.add_slide(slide_layout)
                        title_placeholder = slide.shapes.title
                        subtitle_placeholder = slide.placeholders[1]

                        title_placeholder.text = title
                        period_text = f"Año {period_start.year}" if period_start.year == period_end.year else f"{period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}"
                        subtitle_placeholder.text = f"Período: {period_text}"

                        # Content slide
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title_placeholder = slide.shapes.title
                        content_placeholder = slide.placeholders[1]

                        title_placeholder.text = "Resumen de Actividades"

                        summary = calculate_activity_summary(filtered_forms)
                        content_text = f"""Cursos de Capacitación: {summary['cursos']}
Publicaciones: {summary['publicaciones']}
Eventos Académicos: {summary['eventos']}
Diseños Curriculares: {summary['disenos']}
Experiencias de Movilidad: {summary['movilidades']}
Reconocimientos: {summary['reconocimientos']}
Certificaciones: {summary['certificaciones']}"""

                        content_placeholder.text = content_text

                        buffer = BytesIO()
                        prs.save(buffer)
                        buffer.seek(0)
                        ppt_content = buffer.getvalue()

                        st.download_button(
                            label=f"📄 Exportar como {export_format}",
                            data=ppt_content,
                            file_name=f"{filename_base}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            key=f"direct_ppt_{datetime.now().timestamp()}"
                        )
                    except ImportError:
                        report_content = generate_simple_report(
                            filtered_forms, title, report_type, period_start, period_end)
                        st.download_button(
                            label=f"📄 Exportar como {export_format}",
                            data=report_content,
                            file_name=f"{filename_base}.txt",
                            mime="text/plain",
                            key=f"direct_ppt_txt_{datetime.now().timestamp()}"
                        )

        else:
            st.info("No hay datos disponibles para el período seleccionado.")

    if history_manager:
        with tab2:
            # Report history using new ReportHistory class
            if history_manager:
                history_manager.show_history_interface()
            else:
                st.info("Funcionalidad de historial no disponible")

        with tab3:
            # Storage statistics using new ReportHistory class
            if history_manager:
                st.subheader("📊 Estadísticas de Reportes")
                stats = history_manager.get_statistics()

                col1, col2, col3, col4 = st.columns(4)

                with col1:
                    st.metric("Total Reportes", stats['total_reportes'])

                with col2:
                    st.metric("Última Semana", stats['reportes_ultima_semana'])

                with col3:
                    st.metric("Último Mes", stats['reportes_ultimo_mes'])

                with col4:
                    if stats.get('reporte_mas_reciente'):
                        fecha_reciente = datetime.fromisoformat(
                            stats['reporte_mas_reciente'])
                        st.metric("Más Reciente",
                                  fecha_reciente.strftime('%d/%m/%Y'))
                    else:
                        st.metric("Más Reciente", "N/A")

                # Mostrar tipos más comunes
                if stats['tipos_mas_comunes']:
                    st.subheader("📈 Tipos de Reportes Más Generados")
                    for tipo, count in list(stats['tipos_mas_comunes'].items())[:5]:
                        st.write(f"- **{tipo}**: {count} reportes")
            else:
                st.info("Estadísticas no disponibles")


@st.cache_data(ttl=60)  # Reduced cache time for more frequent updates
def load_report_data():
    """Load data for report generation with caching"""
    db = SessionLocal()
    try:
        crud = FormularioCRUD(db)
        all_forms = crud.get_all_formularios(limit=1000)
        metrics = crud.get_metricas_generales()
        return all_forms, metrics
    finally:
        db.close()


def filter_forms_by_period(forms, start_date, end_date):
    """Filter forms by date period"""
    return [
        f for f in forms
        if f.fecha_envio and start_date <= f.fecha_envio.date() <= end_date
    ]


def get_quarter_dates(quarter, year):
    """Get start and end dates for a quarter"""
    quarter_starts = {
        1: (1, 1), 2: (4, 1), 3: (7, 1), 4: (10, 1)
    }
    quarter_ends = {
        1: (3, 31), 2: (6, 30), 3: (9, 30), 4: (12, 31)
    }

    start_month, start_day = quarter_starts[quarter]
    end_month, end_day = quarter_ends[quarter]

    return (
        date(year, start_month, start_day),
        date(year, end_month, end_day)
    )


def get_quarter_name(quarter):
    """Get quarter name in Spanish"""
    names = {1: 'Primer Trimestre', 2: 'Segundo Trimestre',
             3: 'Tercer Trimestre', 4: 'Cuarto Trimestre'}
    return names[quarter]


def create_preview_dataframe(forms):
    """Create DataFrame for preview - only approved forms with fresh data"""
    data = []
    # Filter only approved forms for preview
    approved_forms = [f for f in forms if f.estado.value == 'APROBADO']

    # Use a single database connection for efficiency
    db = SessionLocal()
    try:
        crud = FormularioCRUD(db)

        for form in approved_forms:
            # Get fresh form with all relationships loaded
            fresh_form = crud.get_formulario(form.id)

            if fresh_form:
                # Count relationships from fresh form
                total_cursos = len(
                    fresh_form.cursos_capacitacion) if fresh_form.cursos_capacitacion else 0
                total_publicaciones = len(
                    fresh_form.publicaciones) if fresh_form.publicaciones else 0
                total_eventos = len(
                    fresh_form.eventos_academicos) if fresh_form.eventos_academicos else 0
                total_disenos = len(
                    fresh_form.diseno_curricular) if fresh_form.diseno_curricular else 0
                total_movilidades = len(
                    fresh_form.movilidad) if fresh_form.movilidad else 0
                total_reconocimientos = len(
                    fresh_form.reconocimientos) if fresh_form.reconocimientos else 0
                total_certificaciones = len(
                    fresh_form.certificaciones) if fresh_form.certificaciones else 0
            else:
                # Fallback to zero counts if form not found
                total_cursos = total_publicaciones = total_eventos = 0
                total_disenos = total_movilidades = total_reconocimientos = total_certificaciones = 0

            data.append({
                'ID': form.id,
                'Docente': form.nombre_completo,
                'Estado': form.estado.value,
                'Fecha': form.fecha_envio.strftime('%Y-%m-%d') if form.fecha_envio else '',
                'Cursos': total_cursos,
                'Publicaciones': total_publicaciones,
                'Eventos': total_eventos,
                'Diseños': total_disenos,
                'Movilidades': total_movilidades,
                'Reconocimientos': total_reconocimientos,
                'Certificaciones': total_certificaciones
            })
    finally:
        db.close()

    return pd.DataFrame(data)


def generate_simple_report(forms, title, report_type, period_start, period_end):
    """Generate a narrative report in the requested format with fresh data"""

    # Calculate basic statistics
    approved_forms = [f for f in forms if f.estado.value == 'APROBADO']

    # Extract detailed activities from approved forms using fresh data
    all_publicaciones = []
    all_cursos = []
    all_eventos = []
    all_disenos = []
    all_movilidades = []
    all_reconocimientos = []
    all_certificaciones = []

    # Use a single database connection for efficiency
    db = SessionLocal()
    try:
        crud = FormularioCRUD(db)

        for form in approved_forms:
            # Get fresh form with all relationships loaded
            fresh_form = crud.get_formulario(form.id)

            if not fresh_form:
                continue
            # Extract publicaciones
            try:
                if fresh_form.publicaciones:
                    for pub in fresh_form.publicaciones:
                        all_publicaciones.append({
                            'titulo': getattr(pub, 'titulo', ''),
                            'autores': getattr(pub, 'autores', ''),
                            'evento_revista': getattr(pub, 'evento_revista', ''),
                            'estatus': getattr(pub, 'estatus', '').value if hasattr(getattr(pub, 'estatus', None), 'value') else str(getattr(pub, 'estatus', ''))
                        })
            except:
                pass

            # Extract cursos
            try:
                if fresh_form.cursos_capacitacion:
                    for curso in fresh_form.cursos_capacitacion:
                        all_cursos.append({
                            'nombre': getattr(curso, 'nombre_curso', ''),
                            'horas': getattr(curso, 'horas', 0),
                            'fecha': getattr(curso, 'fecha', None)
                        })
            except:
                pass

            # Extract eventos
            try:
                if fresh_form.eventos_academicos:
                    for evento in fresh_form.eventos_academicos:
                        all_eventos.append({
                            'nombre': getattr(evento, 'nombre_evento', ''),
                            'tipo': getattr(evento, 'tipo_participacion', '').value if hasattr(getattr(evento, 'tipo_participacion', None), 'value') else str(getattr(evento, 'tipo_participacion', ''))
                        })
            except:
                pass

            # Extract diseños curriculares
            try:
                if fresh_form.diseno_curricular:
                    for diseno in fresh_form.diseno_curricular:
                        all_disenos.append({
                            'nombre': getattr(diseno, 'nombre_curso', ''),
                            'descripcion': getattr(diseno, 'descripcion', '')
                        })
            except:
                pass

            # Extract movilidades
            try:
                if fresh_form.movilidad:
                    for movilidad in fresh_form.movilidad:
                        all_movilidades.append({
                            'descripcion': getattr(movilidad, 'descripcion', ''),
                            'tipo': getattr(movilidad, 'tipo', '').value if hasattr(getattr(movilidad, 'tipo', None), 'value') else str(getattr(movilidad, 'tipo', '')),
                            'fecha': getattr(movilidad, 'fecha', None)
                        })
            except:
                pass

            # Extract reconocimientos
            try:
                if fresh_form.reconocimientos:
                    for reconocimiento in fresh_form.reconocimientos:
                        all_reconocimientos.append({
                            'nombre': getattr(reconocimiento, 'nombre', ''),
                            'tipo': getattr(reconocimiento, 'tipo', '').value if hasattr(getattr(reconocimiento, 'tipo', None), 'value') else str(getattr(reconocimiento, 'tipo', '')),
                            'fecha': getattr(reconocimiento, 'fecha', None)
                        })
            except:
                pass

            # Extract certificaciones
            try:
                if fresh_form.certificaciones:
                    for certificacion in fresh_form.certificaciones:
                        all_certificaciones.append({
                            'nombre': getattr(certificacion, 'nombre', ''),
                            'fecha_obtencion': getattr(certificacion, 'fecha_obtencion', None),
                            'vigente': getattr(certificacion, 'vigente', True)
                        })
            except:
                pass

    finally:
        db.close()

    # Generate narrative report based on type
    if report_type == "annual":
        return generate_annual_narrative_report(
            title, period_start, period_end, approved_forms,
            all_publicaciones, all_cursos, all_eventos, all_disenos,
            all_movilidades, all_reconocimientos, all_certificaciones
        )
    elif report_type == "quarterly":
        return generate_quarterly_narrative_report(
            title, period_start, period_end, approved_forms,
            all_publicaciones, all_cursos, all_eventos, all_disenos,
            all_movilidades, all_reconocimientos, all_certificaciones
        )
    else:
        return generate_data_table_report(
            title, period_start, period_end, approved_forms,
            all_publicaciones, all_cursos, all_eventos, all_disenos,
            all_movilidades, all_reconocimientos, all_certificaciones
        )


def generate_annual_narrative_report(title, period_start, period_end, approved_forms, publicaciones, cursos, eventos, disenos, movilidades, reconocimientos, certificaciones):
    """Generate annual narrative report with examples and detailed descriptions"""

    # Count totals
    total_publicaciones = len(publicaciones)
    total_cursos = len(cursos)
    total_eventos = len(eventos)
    total_disenos = len(disenos)
    total_movilidades = len(movilidades)
    total_reconocimientos = len(reconocimientos)
    total_certificaciones = len(certificaciones)
    total_docentes = len(approved_forms)

    # Get examples for narrative
    pub_examples = []
    if publicaciones:
        for pub in publicaciones[:4]:  # First 4 as examples
            if pub['titulo']:
                tipo = "artículo" if "artículo" in pub['evento_revista'].lower(
                ) else "ponencia" if "ponencia" in pub['evento_revista'].lower() else "publicación"
                pub_examples.append(f"{tipo} *\"{pub['titulo']}\"*")

    curso_examples = []
    if cursos:
        unique_cursos = list(
            {curso['nombre']: curso for curso in cursos if curso['nombre']}.values())
        for curso in unique_cursos[:4]:  # First 4 unique courses
            curso_examples.append(f"*\"{curso['nombre']}\"*")

    evento_examples = []
    if eventos:
        unique_eventos = list(
            {evento['nombre']: evento for evento in eventos if evento['nombre']}.values())
        for evento in unique_eventos[:5]:  # First 5 unique events
            evento_examples.append(f"*{evento['nombre']}*")

    diseno_examples = []
    if disenos:
        unique_disenos = list(
            {diseno['nombre']: diseno for diseno in disenos if diseno['nombre']}.values())
        for diseno in unique_disenos[:4]:  # First 4 unique designs
            diseno_examples.append(f"*{diseno['nombre']}*")

    # Build narrative report
    report_lines = [
        f"# {title}",
        "",
        f"**Período:** Año {period_start.year}" if period_start.year == period_end.year else f"**Período:** {period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}",
        f"**Fecha de generación:** {datetime.now().strftime('%d de %B de %Y')}",
        "",
        "## Resumen de Actividades Académicas",
        "",
        f"En el Departamento se realizaron los siguientes productos durante el período {period_start.year}:",
        ""
    ]

    # Publications section
    if total_publicaciones > 0:
        pub_text = f"{total_publicaciones} trabajos de publicación"
        if pub_examples:
            ejemplos = ", ".join(pub_examples[:3])
            if len(pub_examples) > 3:
                ejemplos += f" y otros {len(pub_examples) - 3} trabajos más"
            pub_text += f", entre ellos: {ejemplos}"
        report_lines.append(f"> {pub_text}.")
        report_lines.append("")

    # Training courses section
    if total_cursos > 0:
        curso_text = f"{total_docentes} docentes se capacitaron en {total_cursos} cursos"
        if curso_examples:
            ejemplos = ", ".join(curso_examples[:4])
            curso_text += f" como {ejemplos}"
        report_lines.append(f"> {curso_text}.")
        report_lines.append("")

    # Curriculum design section
    if total_disenos > 0:
        diseno_text = f"{total_disenos} productos de Diseño Curricular liberados"
        if diseno_examples:
            ejemplos = ", ".join(diseno_examples[:4])
            diseno_text += f", entre ellos cursos como {ejemplos}"
        report_lines.append(f"> {diseno_text}.")
        report_lines.append("")

    # Academic events section
    if total_eventos > 0:
        evento_text = f"{total_eventos} eventos académicos organizados"
        if evento_examples:
            ejemplos = ", ".join(evento_examples[:5])
            evento_text += f", tales como {ejemplos}"
        report_lines.append(f"> {evento_text}.")
        report_lines.append("")

    # Mobility experiences section
    if total_movilidades > 0:
        movilidad_text = f"{total_movilidades} experiencias de movilidad académica realizadas"
        report_lines.append(f"> {movilidad_text}.")
        report_lines.append("")

    # Recognitions section
    if total_reconocimientos > 0:
        reconocimiento_text = f"{total_reconocimientos} reconocimientos y distinciones obtenidos"
        report_lines.append(f"> {reconocimiento_text}.")
        report_lines.append("")

    # Certifications section
    if total_certificaciones > 0:
        certificacion_text = f"{total_certificaciones} certificaciones profesionales adquiridas"
        report_lines.append(f"> {certificacion_text}.")
        report_lines.append("")

    # Summary statistics
    report_lines.extend([
        "## Estadísticas Generales",
        "",
        f"- **Total de docentes participantes:** {total_docentes}",
        f"- **Publicaciones generadas:** {total_publicaciones}",
        f"- **Cursos de capacitación:** {total_cursos}",
        f"- **Eventos académicos:** {total_eventos}",
        f"- **Diseños curriculares:** {total_disenos}",
        f"- **Experiencias de movilidad:** {total_movilidades}",
        f"- **Reconocimientos:** {total_reconocimientos}",
        f"- **Certificaciones:** {total_certificaciones}",
        "",
        "---",
        "",
        f"*Reporte generado automáticamente por el Sistema de Reportes Docentes*",
        f"*Fecha: {datetime.now().strftime('%d de %B de %Y a las %H:%M')}*"
    ])

    return "\n".join(report_lines)


def generate_quarterly_narrative_report(title, period_start, period_end, approved_forms, publicaciones, cursos, eventos, disenos, movilidades, reconocimientos, certificaciones):
    """Generate quarterly narrative report with hard data and brief examples"""

    # Count totals
    total_publicaciones = len(publicaciones)
    total_cursos = len(cursos)
    total_eventos = len(eventos)
    total_disenos = len(disenos)
    total_movilidades = len(movilidades)
    total_reconocimientos = len(reconocimientos)
    total_certificaciones = len(certificaciones)
    total_docentes = len(approved_forms)

    # Get brief examples
    pub_examples = [pub['titulo']
                    for pub in publicaciones[:2] if pub['titulo']]
    curso_examples = list({curso['nombre']
                          for curso in cursos if curso['nombre']})[:3]
    evento_examples = list({evento['nombre']
                           for evento in eventos if evento['nombre']})[:3]
    diseno_examples = list({diseno['nombre']
                           for diseno in disenos if diseno['nombre']})[:3]

    # Determine quarter name
    quarter_num = ((period_start.month - 1) // 3) + 1
    quarter_names = {1: "1er", 2: "2do", 3: "3er", 4: "4to"}
    quarter_name = quarter_names.get(quarter_num, str(quarter_num))

    # Build report
    report_lines = [
        f"# {title}",
        "",
        f"**Período:** {quarter_name} {period_start.year}",
        f"**Fecha de generación:** {datetime.now().strftime('%d de %B de %Y')}",
        "",
        f"**{quarter_name} {period_start.year} ({period_start.strftime('%B').lower()}-{period_end.strftime('%B').lower()}):**",
        ""
    ]

    # Publications
    if total_publicaciones > 0:
        pub_text = f"- {total_publicaciones} artículos publicados en revistas indexadas"
        if pub_examples:
            ejemplos = ", ".join([f"*{pub}*" for pub in pub_examples[:2]])
            pub_text += f" (ejemplo: {ejemplos})"
        report_lines.append(f"> {pub_text}.")

    # Training
    if total_cursos > 0:
        curso_text = f"- {total_docentes} docentes capacitados en cursos"
        if curso_examples:
            ejemplos = ", ".join(
                [f"*{curso}*" for curso in curso_examples[:2]])
            curso_text += f" (ejemplo: {ejemplos})"
        report_lines.append(f"> {curso_text}.")

    # Curriculum design
    if total_disenos > 0:
        diseno_text = f"- {total_disenos} diseños curriculares liberados"
        if diseno_examples:
            ejemplos = ", ".join(
                [f"*{diseno}*" for diseno in diseno_examples[:3]])
            diseno_text += f" (ejemplo: {ejemplos})"
        report_lines.append(f"> {diseno_text}.")

    # Events
    if total_eventos > 0:
        evento_text = f"- {total_eventos} eventos académicos organizados"
        if evento_examples:
            ejemplos = ", ".join(
                [f"*{evento}*" for evento in evento_examples[:3]])
            evento_text += f" ({ejemplos})"
        report_lines.append(f"> {evento_text}.")

    report_lines.extend([
        "",
        "---",
        "",
        f"*Reporte trimestral generado automáticamente*",
        f"*{datetime.now().strftime('%d de %B de %Y')}*"
    ])

    return "\n".join(report_lines)


def generate_data_table_report(title, period_start, period_end, approved_forms, publicaciones, cursos, eventos, disenos, movilidades, reconocimientos, certificaciones):
    """Generate data table report with structured information"""

    report_lines = [
        f"# {title}",
        "",
        f"**Período:** Año {period_start.year}" if period_start.year == period_end.year else f"**Período:** {period_start.strftime('%Y-%m-%d')} - {period_end.strftime('%Y-%m-%d')}",
        f"**Fecha de generación:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
        "",
        "## Resumen de Datos",
        "",
        f"| Categoría | Cantidad |",
        f"|-----------|----------|",
        f"| Docentes participantes | {len(approved_forms)} |",
        f"| Publicaciones | {len(publicaciones)} |",
        f"| Cursos de capacitación | {len(cursos)} |",
        f"| Eventos académicos | {len(eventos)} |",
        f"| Diseños curriculares | {len(disenos)} |",
        f"| Experiencias de movilidad | {len(movilidades)} |",
        f"| Reconocimientos | {len(reconocimientos)} |",
        f"| Certificaciones | {len(certificaciones)} |",
        "",
        "## Detalle de Publicaciones",
        ""
    ]

    if publicaciones:
        report_lines.append("| Título | Autores | Revista/Evento |")
        report_lines.append("|--------|---------|----------------|")
        for pub in publicaciones[:10]:  # First 10
            titulo = pub['titulo'][:50] + \
                "..." if len(pub['titulo']) > 50 else pub['titulo']
            autores = pub['autores'][:30] + \
                "..." if len(pub['autores']) > 30 else pub['autores']
            evento = pub['evento_revista'][:30] + \
                "..." if len(pub['evento_revista']
                             ) > 30 else pub['evento_revista']
            report_lines.append(f"| {titulo} | {autores} | {evento} |")
    else:
        report_lines.append(
            "*No hay publicaciones registradas para este período.*")

    report_lines.extend([
        "",
        "## Detalle de Cursos de Capacitación",
        ""
    ])

    if cursos:
        unique_cursos = {}
        for curso in cursos:
            if curso['nombre'] in unique_cursos:
                unique_cursos[curso['nombre']]['count'] += 1
                unique_cursos[curso['nombre']]['horas'] += curso['horas']
            else:
                unique_cursos[curso['nombre']] = {
                    'count': 1,
                    'horas': curso['horas']
                }

        report_lines.append("| Curso | Participantes | Horas |")
        report_lines.append("|-------|---------------|-------|")
        for nombre, data in unique_cursos.items():
            report_lines.append(
                f"| {nombre} | {data['count']} | {data['horas']} |")
    else:
        report_lines.append("*No hay cursos registrados para este período.*")

    report_lines.extend([
        "",
        "---",
        "",
        f"*Reporte de datos generado automáticamente*",
        f"*{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*"
    ])

    return "\n".join(report_lines)


def calculate_activity_summary(forms):
    """Calculate activity summary for approved forms with fresh data"""
    approved_forms = [f for f in forms if f.estado.value == 'APROBADO']

    total_cursos = 0
    total_publicaciones = 0
    total_eventos = 0
    total_disenos = 0
    total_movilidades = 0
    total_reconocimientos = 0
    total_certificaciones = 0

    # Use a single database connection for efficiency
    db = SessionLocal()
    try:
        crud = FormularioCRUD(db)

        for f in approved_forms:
            # Get fresh form with all relationships loaded
            fresh_form = crud.get_formulario(f.id)

            if fresh_form:
                # Count relationships from fresh form
                total_cursos += len(
                    fresh_form.cursos_capacitacion) if fresh_form.cursos_capacitacion else 0
                total_publicaciones += len(
                    fresh_form.publicaciones) if fresh_form.publicaciones else 0
                total_eventos += len(
                    fresh_form.eventos_academicos) if fresh_form.eventos_academicos else 0
                total_disenos += len(
                    fresh_form.diseno_curricular) if fresh_form.diseno_curricular else 0
                total_movilidades += len(
                    fresh_form.movilidad) if fresh_form.movilidad else 0
                total_reconocimientos += len(
                    fresh_form.reconocimientos) if fresh_form.reconocimientos else 0
                total_certificaciones += len(
                    fresh_form.certificaciones) if fresh_form.certificaciones else 0
    finally:
        db.close()

    return {
        'cursos': total_cursos,
        'publicaciones': total_publicaciones,
        'eventos': total_eventos,
        'disenos': total_disenos,
        'movilidades': total_movilidades,
        'reconocimientos': total_reconocimientos,
        'certificaciones': total_certificaciones
    }


def export_report_basic(forms, report_type, period_start, period_end, export_format, title):
    """Export report in native formats (PDF, Excel, PowerPoint)"""

    def generate_pdf_content(forms, title, report_type, period_start, period_end):
        """Generate PDF content with same design as Markdown report"""
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib import colors
            from io import BytesIO

            buffer = BytesIO()
            doc = SimpleDocTemplate(
                buffer, pagesize=A4, topMargin=50, bottomMargin=50)
            styles = getSampleStyleSheet()
            story = []

            # Custom styles
            title_style = ParagraphStyle('CustomTitle', parent=styles['Heading1'],
                                         fontSize=20, spaceAfter=30, alignment=1, textColor=colors.HexColor('#1f77b4'))

            heading_style = ParagraphStyle('CustomHeading', parent=styles['Heading2'],
                                           fontSize=16, spaceAfter=20, textColor=colors.HexColor('#2e7d32'))

            # Generate the same content as Markdown report
            report_content = generate_simple_report(
                forms, title, report_type, period_start, period_end)

            # Parse and format content
            lines = report_content.split('\n')
            for line in lines:
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 6))
                elif line.startswith('# '):
                    story.append(Paragraph(line[2:], title_style))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], heading_style))
                elif line.startswith('> '):
                    # Highlight boxes for activities
                    content = line[2:]
                    highlight_style = ParagraphStyle('Highlight', parent=styles['Normal'],
                                                     leftIndent=20, rightIndent=20,
                                                     backColor=colors.HexColor(
                                                         '#f0f8ff'),
                                                     borderColor=colors.HexColor(
                                                         '#1f77b4'),
                                                     borderWidth=1, borderPadding=10)
                    story.append(Paragraph(content, highlight_style))
                elif line.startswith('- **'):
                    # Statistics lines
                    story.append(Paragraph(line, styles['Normal']))
                elif line.startswith('*') and line.endswith('*'):
                    # Italic footer text
                    italic_style = ParagraphStyle('Italic', parent=styles['Normal'],
                                                  fontName='Helvetica-Oblique', fontSize=10,
                                                  alignment=1, textColor=colors.grey)
                    story.append(Paragraph(line[1:-1], italic_style))
                else:
                    story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 6))

            doc.build(story)
            buffer.seek(0)
            return buffer.getvalue()
        except ImportError:
            content = generate_simple_report(
                forms, title, report_type, period_start, period_end)
            return content.encode('utf-8')

    def generate_excel_content(forms, title, report_type, period_start, period_end):
        """Generate Excel content with same design as Markdown report"""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
            from io import BytesIO

            wb = Workbook()

            # Sheet 1: Report Content (same as Markdown)
            ws1 = wb.active
            ws1.title = "Reporte Narrativo"

            # Generate the same content as Markdown report
            report_content = generate_simple_report(
                forms, title, report_type, period_start, period_end)

            # Add report content to Excel (simplified without merging)
            row = 1
            lines = report_content.split('\n')
            for line in lines:
                line = line.strip()
                if line:
                    cell = ws1.cell(row=row, column=1, value=line)

                    if line.startswith('# '):
                        # Main title
                        cell.value = line[2:]
                        cell.font = Font(size=18, bold=True, color='1f77b4')
                    elif line.startswith('## '):
                        # Section headers
                        cell.value = line[3:]
                        cell.font = Font(size=14, bold=True, color='2e7d32')
                    elif line.startswith('> '):
                        # Activity highlights
                        cell.value = line[2:]
                        cell.fill = PatternFill(
                            start_color='f0f8ff', end_color='f0f8ff', fill_type='solid')
                    elif line.startswith('- **'):
                        # Statistics
                        cell.font = Font(bold=True)

                    row += 1
                else:
                    row += 1  # Empty line

            # Sheet 2: Data Table
            ws2 = wb.create_sheet("Datos Detallados")

            # Headers
            headers = ['ID', 'Docente', 'Estado', 'Fecha', 'Cursos', 'Publicaciones',
                       'Eventos', 'Diseños', 'Movilidades', 'Reconocimientos', 'Certificaciones']
            for col, header in enumerate(headers, 1):
                cell = ws2.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True, color='ffffff')
                cell.fill = PatternFill(
                    start_color='1f77b4', end_color='1f77b4', fill_type='solid')
                cell.alignment = Alignment(horizontal='center')

            # Data rows
            approved_forms = [f for f in forms if f.estado.value == 'APROBADO']
            db = SessionLocal()
            try:
                crud = FormularioCRUD(db)
                for row, form in enumerate(approved_forms, 2):
                    fresh_form = crud.get_formulario(form.id)
                    if fresh_form:
                        ws2.cell(row=row, column=1, value=form.id)
                        ws2.cell(row=row, column=2, value=form.nombre_completo)
                        ws2.cell(row=row, column=3, value=form.estado.value)
                        ws2.cell(row=row, column=4, value=form.fecha_envio.strftime(
                            '%Y-%m-%d') if form.fecha_envio else '')
                        ws2.cell(row=row, column=5, value=len(
                            fresh_form.cursos_capacitacion) if fresh_form.cursos_capacitacion else 0)
                        ws2.cell(row=row, column=6, value=len(
                            fresh_form.publicaciones) if fresh_form.publicaciones else 0)
                        ws2.cell(row=row, column=7, value=len(
                            fresh_form.eventos_academicos) if fresh_form.eventos_academicos else 0)
                        ws2.cell(row=row, column=8, value=len(
                            fresh_form.diseno_curricular) if fresh_form.diseno_curricular else 0)
                        ws2.cell(row=row, column=9, value=len(
                            fresh_form.movilidad) if fresh_form.movilidad else 0)
                        ws2.cell(row=row, column=10, value=len(
                            fresh_form.reconocimientos) if fresh_form.reconocimientos else 0)
                        ws2.cell(row=row, column=11, value=len(
                            fresh_form.certificaciones) if fresh_form.certificaciones else 0)
            finally:
                db.close()

            # Set fixed column widths to avoid merged cell issues
            # Wide column for report content
            ws1.column_dimensions['A'].width = 80

            # Auto-adjust column widths for data sheet only
            for col_num in range(1, len(headers) + 1):
                column_letter = ws2.cell(row=1, column=col_num).column_letter
                if col_num == 2:  # Docente column
                    ws2.column_dimensions[column_letter].width = 25
                elif col_num == 4:  # Fecha column
                    ws2.column_dimensions[column_letter].width = 12
                else:
                    ws2.column_dimensions[column_letter].width = 10

            buffer = BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        except ImportError:
            content = generate_simple_report(
                forms, title, report_type, period_start, period_end)
            return content.encode('utf-8')

    def generate_powerpoint_content(forms, title, report_type, period_start, period_end):
        """Generate PowerPoint content with same design as Markdown report"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from io import BytesIO

            prs = Presentation()

            # Generate the same content as Markdown report
            report_content = generate_simple_report(
                forms, title, report_type, period_start, period_end)

            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title_placeholder = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]

            title_placeholder.text = title
            subtitle_placeholder.text = f"Período: Año {period_start.year}" if period_start.year == period_end.year else f"Período: {period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}"

            # Parse report content and create slides
            lines = report_content.split('\n')
            current_slide = None
            current_content = []

            for line in lines:
                line = line.strip()
                if line.startswith('## ') and not line.startswith('## Estadísticas'):
                    # New section - create new slide
                    if current_slide and current_content:
                        # Add content to previous slide
                        if len(current_slide.placeholders) > 1:
                            content_placeholder = current_slide.placeholders[1]
                            content_placeholder.text = '\n'.join(
                                current_content)

                    # Create new slide
                    slide_layout = prs.slide_layouts[1]
                    current_slide = prs.slides.add_slide(slide_layout)
                    title_placeholder = current_slide.shapes.title
                    title_placeholder.text = line[3:]  # Remove "## "
                    current_content = []

                elif line.startswith('> '):
                    # Activity content
                    current_content.append(f"* {line[2:]}")

                elif line.startswith('## Estadísticas'):
                    # Statistics slide
                    if current_slide and current_content:
                        if len(current_slide.placeholders) > 1:
                            content_placeholder = current_slide.placeholders[1]
                            content_placeholder.text = '\n'.join(
                                current_content)

                    # Create statistics slide
                    slide_layout = prs.slide_layouts[1]
                    stats_slide = prs.slides.add_slide(slide_layout)
                    title_placeholder = stats_slide.shapes.title
                    title_placeholder.text = "Estadísticas Generales"

                    # Collect statistics
                    stats_content = []
                    continue

                elif line.startswith('- **') and 'current_slide' in locals():
                    # Statistics line
                    stats_content.append(line[2:])  # Remove "- "

            # Add final content
            if current_slide and current_content:
                if len(current_slide.placeholders) > 1:
                    content_placeholder = current_slide.placeholders[1]
                    content_placeholder.text = '\n'.join(current_content)

            # Add statistics to the last slide if exists
            if 'stats_slide' in locals() and 'stats_content' in locals():
                if len(stats_slide.placeholders) > 1:
                    content_placeholder = stats_slide.placeholders[1]
                    content_placeholder.text = '\n'.join(stats_content)

            # Style all slides
            for slide in prs.slides:
                if slide.shapes.title:
                    title_shape = slide.shapes.title
                    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(
                        31, 119, 180)  # Blue
                    title_shape.text_frame.paragraphs[0].font.size = Pt(28)

                # Style content if exists
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph.font.size = Pt(18)
                            paragraph.space_after = Pt(12)

            buffer = BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
        except ImportError:
            content = generate_simple_report(
                forms, title, report_type, period_start, period_end)
            return content.encode('utf-8')

    try:
        with st.spinner(f"Generando reporte en formato {export_format.upper()}..."):

            if export_format == "pdf":
                pdf_content = generate_pdf_content(
                    forms, title, report_type, period_start, period_end)

                filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pdf"

                # Create hidden download button and auto-click it
                import base64
                b64_pdf = base64.b64encode(pdf_content).decode()

                # Create download button that will be auto-clicked
                st.download_button(
                    label="📥 Descargar PDF",
                    data=pdf_content,
                    file_name=filename,
                    mime="application/pdf",
                    key=f"pdf_download_{datetime.now().timestamp()}",
                    help="Descarga automática iniciada"
                )

                # Auto-click the download button using JavaScript
                st.markdown(f"""
                <script>
                setTimeout(function() {{
                    // Find the download button and click it
                    const buttons = document.querySelectorAll('button[kind="primary"]');
                    for (let button of buttons) {{
                        if (button.textContent.includes('Descargar PDF')) {{
                            button.click();
                            break;
                        }}
                    }}
                }}, 500);
                </script>
                """, unsafe_allow_html=True)

                st.success(f"✅ PDF generado y descargado: {filename}")

            elif export_format == "excel":
                excel_content = generate_excel_content(
                    forms, title, report_type, period_start, period_end)

                filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.xlsx"

                # Create download button that will be auto-clicked
                st.download_button(
                    label="📥 Descargar Excel",
                    data=excel_content,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    key=f"excel_download_{datetime.now().timestamp()}",
                    help="Descarga automática iniciada"
                )

                # Auto-click the download button using JavaScript
                st.markdown(f"""
                <script>
                setTimeout(function() {{
                    // Find the download button and click it
                    const buttons = document.querySelectorAll('button[kind="primary"]');
                    for (let button of buttons) {{
                        if (button.textContent.includes('Descargar Excel')) {{
                            button.click();
                            break;
                        }}
                    }}
                }}, 500);
                </script>
                """, unsafe_allow_html=True)

                st.success(f"✅ Excel generado y descargado: {filename}")

            elif export_format == "powerpoint":
                ppt_content = generate_powerpoint_content(
                    forms, title, report_type, period_start, period_end)

                filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"

                # Create download button that will be auto-clicked
                st.download_button(
                    label="📥 Descargar PowerPoint",
                    data=ppt_content,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"ppt_download_{datetime.now().timestamp()}",
                    help="Descarga automática iniciada"
                )

                # Auto-click the download button using JavaScript
                st.markdown(f"""
                <script>
                setTimeout(function() {{
                    // Find the download button and click it
                    const buttons = document.querySelectorAll('button[kind="primary"]');
                    for (let button of buttons) {{
                        if (button.textContent.includes('Descargar PowerPoint')) {{
                            button.click();
                            break;
                        }}
                    }}
                }}, 500);
                </script>
                """, unsafe_allow_html=True)

                st.success(f"✅ PowerPoint generado y descargado: {filename}")

            st.success(
                f"✅ Reporte en formato {export_format.upper()} generado exitosamente!")

            # Show export details
            with st.expander("ℹ️ Detalles de Exportación"):
                st.write(f"**Formato:** {export_format.upper()}")
                st.write(f"**Título:** {title}")
                st.write(f"**Período:** {period_start} - {period_end}")
                st.write(f"**Formularios procesados:** {len(forms)}")
                st.write(
                    f"**Fecha de generación:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

                if export_format == "excel":
                    st.write("**Hojas incluidas:** Formularios, Resumen")
                elif export_format == "pdf":
                    st.write("**Formato:** HTML optimizado para impresión")
                elif export_format == "powerpoint":
                    st.write("**Diapositivas:** Portada, Resumen, Actividades")

    except Exception as e:
        st.error(f"❌ Error al exportar reporte: {str(e)}")


def generate_simple_report_display(forms, report_type, period_start, period_end, title):
    """Generate and display a simple report without advanced components"""

    try:
        with st.spinner("Generando reporte..."):
            # Generate simple report content
            report_content = generate_simple_report(
                forms, title, report_type, period_start, period_end)

            st.success("✅ Reporte generado exitosamente!")

            # Display report
            st.subheader("📄 Reporte Generado")

            # Add download button
            st.download_button(
                label="📥 Descargar Reporte",
                data=report_content,
                file_name=f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.md",
                mime="text/markdown",
                key="download_simple_report"
            )

            # Display content
            st.markdown(report_content)

            # Show generation details
            with st.expander("ℹ️ Detalles de Generación"):
                st.write(f"**Tipo:** {report_type}")
                st.write(f"**Período:** {period_start} - {period_end}")
                st.write(f"**Formularios procesados:** {len(forms)}")
                st.write(
                    f"**Fecha de generación:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    except Exception as e:
        st.error(f"❌ Error al generar reporte: {str(e)}")


def generate_and_display_report(report_generator, history_manager, forms,
                                report_type, period_start, period_end,
                                include_trends, include_highlights, report_tone, custom_title):
    """Generate report and display results"""

    try:
        with st.spinner("Generando reporte con técnicas de NLG..."):

            # Generate title
            if custom_title:
                title = custom_title
            else:
                if report_type == "annual":
                    title = f"Reporte Anual de Actividades Docentes {period_start.year}"
                elif report_type == "quarterly":
                    quarter = ((period_start.month - 1) // 3) + 1
                    title = f"Resumen Trimestral Q{quarter} {period_start.year}"
                else:
                    title = f"Reporte de Datos {period_start.strftime('%B %Y')}"

            # Generate simple report content
            report_content = generate_simple_report(
                forms, title, report_type, period_start, period_end)

            # Customize tone if needed
            if report_tone != "professional":
                report_content = report_generator.nlg_engine.customize_report_tone(
                    report_content, report_tone
                )

            # Save to history using new ReportHistory class
            report_id = history_manager.add_report({
                'tipo': report_type,
                'formato': 'markdown',
                'total_registros': len(forms),
                'nombre_archivo': f"{title.replace(' ', '_')}.md",
                'descripcion': title,
                'filtros_aplicados': {
                    'periodo_inicio': period_start.isoformat(),
                    'periodo_fin': period_end.isoformat(),
                    'incluye_tendencias': include_trends,
                    'incluye_destacados': include_highlights,
                    'tono': report_tone
                }
            })

            st.success(f"✅ Reporte generado exitosamente! ID: {report_id}")

            # Display report
            st.subheader("📄 Reporte Generado")

            # Add download button
            st.download_button(
                label="📥 Descargar Reporte",
                data=report_content,
                file_name=f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.md",
                mime="text/markdown",
                key="download_report"
            )

            # Display content
            st.markdown(report_content)

            # Show generation details
            with st.expander("ℹ️ Detalles de Generación"):
                st.write(f"**ID del reporte:** {report_id}")
                st.write(f"**Tipo:** {report_type}")
                st.write(f"**Período:** {period_start} - {period_end}")
                st.write(f"**Formularios procesados:** {len(forms)}")
                st.write(f"**Tono:** {report_tone}")
                st.write(
                    f"**Fecha de generación:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    except Exception as e:
        st.error(f"❌ Error al generar reporte: {str(e)}")
        st.exception(e)


def show_report_history(history_manager):
    """Show report history interface"""

    st.subheader("📚 Historial de Reportes")

    # Get history
    history = history_manager.get_report_history(limit=50)

    if not history:
        st.info("No hay reportes generados aún.")
        return

    # Filter options
    col1, col2 = st.columns(2)

    with col1:
        type_filter = st.selectbox(
            "Filtrar por tipo:",
            ["Todos", "annual", "quarterly", "data_table"],
            format_func=lambda x: {
                "Todos": "Todos los tipos",
                "annual": "Reportes Anuales",
                "quarterly": "Reportes Trimestrales",
                "data_table": "Reportes de Datos"
            }.get(x, x)
        )

    with col2:
        search_query = st.text_input(
            "Buscar reportes:", placeholder="Título, resumen...")

    # Apply filters
    filtered_history = history
    if type_filter != "Todos":
        filtered_history = [
            r for r in filtered_history if r['report_type'] == type_filter]

    if search_query:
        filtered_history = history_manager.search_reports(search_query)

    # Display reports
    for report in filtered_history[:20]:  # Show first 20
        with st.expander(f"📄 {report['title']} ({report['generation_date'][:10]})"):

            col1, col2 = st.columns([3, 1])

            with col1:
                st.write(f"**Tipo:** {report['report_type']}")
                st.write(
                    f"**Período:** {report['period_start'][:10]} - {report['period_end'][:10]}")
                st.write(f"**Resumen:** {report['summary']}")
                st.write(f"**Tamaño:** {report['file_size']} bytes")

            with col2:
                # Action buttons
                if st.button("👁️ Ver", key=f"view_{report['id']}"):
                    content = history_manager.get_report_content(report['id'])
                    if content:
                        st.markdown("---")
                        st.markdown(content)

                if st.button("📥 Descargar", key=f"download_{report['id']}"):
                    content = history_manager.get_report_content(report['id'])
                    if content:
                        st.download_button(
                            label="Descargar archivo",
                            data=content,
                            file_name=f"{report['title'].replace(' ', '_')}.md",
                            mime="text/markdown",
                            key=f"dl_{report['id']}"
                        )

                if st.button("🗑️ Eliminar", key=f"delete_{report['id']}"):
                    if history_manager.delete_report(report['id']):
                        st.success("Reporte eliminado")
                        st.rerun()


def show_storage_statistics(history_manager):
    """Show storage statistics"""

    st.subheader("📊 Estadísticas de Almacenamiento")

    stats = history_manager.get_storage_statistics()

    # Main metrics
    col1, col2, col3, col4 = st.columns(4)

    with col1:
        st.metric("Total Reportes", stats['total_reports'])
    with col2:
        st.metric("Tamaño Total", f"{stats['total_size_mb']} MB")
    with col3:
        st.metric("Reportes Recientes", stats['recent_reports_count'])
    with col4:
        if stats['total_reports'] > 0:
            avg_size = stats['total_size_mb'] / stats['total_reports']
            st.metric("Tamaño Promedio", f"{avg_size:.2f} MB")

    # Reports by type
    if stats['reports_by_type']:
        st.subheader("Distribución por Tipo")

        type_df = pd.DataFrame([
            {'Tipo': tipo, 'Cantidad': cantidad}
            for tipo, cantidad in stats['reports_by_type'].items()
        ])

        st.bar_chart(type_df.set_index('Tipo'))

    # Timeline
    if stats['oldest_report'] and stats['newest_report']:
        st.subheader("Línea de Tiempo")
        st.write(f"**Primer reporte:** {stats['oldest_report'][:10]}")
        st.write(f"**Último reporte:** {stats['newest_report'][:10]}")

    # Cleanup options
    st.subheader("🧹 Mantenimiento")

    col1, col2 = st.columns(2)

    with col1:
        days_to_keep = st.number_input(
            "Días a conservar:",
            min_value=30,
            max_value=3650,
            value=365
        )

        if st.button("Limpiar Reportes Antiguos"):
            deleted_count = history_manager.cleanup_old_reports(days_to_keep)
            st.success(f"Se eliminaron {deleted_count} reportes antiguos")

    with col2:
        if st.button("Exportar Metadatos"):
            metadata = history_manager.export_history_metadata()
            st.download_button(
                label="Descargar Metadatos",
                data=metadata,
                file_name=f"report_metadata_{datetime.now().strftime('%Y%m%d')}.json",
                mime="application/json"
            )


if __name__ == "__main__":
    show_report_generation_page()


def export_multiformat_report(report_generator, forms, report_type,
                              period_start, period_end, export_format,
                              custom_title, include_charts, include_metadata):
    """Export report in multiple formats (PDF, Excel, PowerPoint)"""

    try:
        with st.spinner(f"Generando reporte en formato {export_format.upper()}..."):

            # Generate title
            if custom_title:
                title = custom_title
            else:
                if report_type == "annual":
                    title = f"Reporte Anual de Actividades Docentes {period_start.year}"
                elif report_type == "quarterly":
                    quarter = ((period_start.month - 1) // 3) + 1
                    title = f"Resumen Trimestral Q{quarter} {period_start.year}"
                else:
                    title = f"Reporte de Datos {period_start.strftime('%B %Y')}"

            # Export report
            if export_format == 'pdf':
                export_format_name = 'PDF'
                mime_type = 'application/pdf'
                file_extension = 'pdf'
            elif export_format == 'excel':
                export_format_name = 'Excel'
                mime_type = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                file_extension = 'xlsx'
            elif export_format == 'powerpoint':
                export_format_name = 'PowerPoint'
                mime_type = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
                file_extension = 'pptx'
            else:
                st.error(f"Formato no soportado: {export_format}")
                return

            # Generate the export
            exported_content = report_generator.export_report_to_format(
                forms=forms,
                report_type=report_type,
                period_start=period_start,
                period_end=period_end,
                export_format=export_format,
                title=title,
                include_charts=include_charts,
                include_metadata=include_metadata
            )

            # Create filename
            safe_title = title.replace(' ', '_').replace('/', '_')
            filename = f"{safe_title}_{datetime.now().strftime('%Y%m%d')}.{file_extension}"

            # Provide download
            st.success(
                f"✅ Reporte en formato {export_format_name} generado exitosamente!")

            st.download_button(
                label=f"📥 Descargar {export_format_name}",
                data=exported_content.getvalue(),
                file_name=filename,
                mime=mime_type,
                key=f"download_{export_format}"
            )

            # Show export details
            with st.expander("ℹ️ Detalles de Exportación"):
                st.write(f"**Formato:** {export_format_name}")
                st.write(f"**Título:** {title}")
                st.write(f"**Período:** {period_start} - {period_end}")
                st.write(f"**Formularios procesados:** {len(forms)}")
                st.write(
                    f"**Gráficos incluidos:** {'Sí' if include_charts else 'No'}")
                st.write(
                    f"**Metadatos incluidos:** {'Sí' if include_metadata else 'No'}")
                st.write(
                    f"**Tamaño del archivo:** {len(exported_content.getvalue())} bytes")
                st.write(
                    f"**Fecha de generación:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

                # Format-specific information
                if export_format == 'pdf':
                    st.write("**Características PDF:**")
                    st.write("- Formato profesional con estilos personalizados")
                    st.write("- Gráficos integrados usando ReportLab")
                    st.write("- Tablas formateadas con colores institucionales")
                    st.write("- Metadatos incluidos en sección dedicada")

                elif export_format == 'excel':
                    st.write("**Características Excel:**")
                    st.write("- Múltiples hojas con datos organizados")
                    st.write("- Gráficos interactivos en hoja dedicada")
                    st.write("- Formato profesional con colores y estilos")
                    st.write("- Tablas con filtros y formato automático")

                elif export_format == 'powerpoint':
                    st.write("**Características PowerPoint:**")
                    st.write("- Presentación con diapositivas estructuradas")
                    st.write("- Gráficos integrados en diapositivas")
                    st.write("- Diseño profesional con colores institucionales")
                    st.write("- Contenido organizado por secciones")

    except Exception as e:
        st.error(
            f"❌ Error al generar reporte en formato {export_format}: {str(e)}")
        st.exception(e)  # Show full traceback for debugging


def generate_pdf_report(forms, title, report_type, period_start, period_end):
    """Generate PDF report using reportlab"""
    try:
        from reportlab.lib.pagesizes import letter, A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from io import BytesIO

        # Create PDF buffer
        buffer = BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1  # Center
        )
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 12))

        # Period info
        period_text = f"Período: Año {period_start.year}" if period_start.year == period_end.year else f"Período: {period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}"
        story.append(Paragraph(period_text, styles['Normal']))
        story.append(Spacer(1, 12))

        # Generate report content
        report_content = generate_simple_report(
            forms, title, report_type, period_start, period_end)

        # Convert markdown to paragraphs
        lines = report_content.split('\n')
        for line in lines:
            if line.strip():
                if line.startswith('# '):
                    story.append(Paragraph(line[2:], styles['Heading1']))
                elif line.startswith('## '):
                    story.append(Paragraph(line[3:], styles['Heading2']))
                elif line.startswith('> '):
                    story.append(Paragraph(line[2:], styles['Normal']))
                else:
                    story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 6))

        # Build PDF
        doc.build(story)
        buffer.seek(0)
        return buffer.getvalue()

    except ImportError:
        # Fallback to simple text-based PDF
        content = generate_simple_report(
            forms, title, report_type, period_start, period_end)
        return content.encode('utf-8')


def generate_excel_report(forms, title, report_type, period_start, period_end):
    """Generate Excel report using openpyxl"""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment, PatternFill
        from io import BytesIO

        # Create workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Reporte Docentes"

        # Header
        ws['A1'] = title
        ws['A1'].font = Font(size=16, bold=True)
        ws['A1'].alignment = Alignment(horizontal='center')
        ws.merge_cells('A1:H1')

        ws['A2'] = f"Período: Año {period_start.year}" if period_start.year == period_end.year else f"Período: {period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}"
        ws['A2'].font = Font(size=12)
        ws.merge_cells('A2:H2')

        # Data headers
        headers = ['ID', 'Docente', 'Estado', 'Fecha', 'Cursos',
                   'Publicaciones', 'Eventos', 'Certificaciones']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=4, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(
                start_color='CCCCCC', end_color='CCCCCC', fill_type='solid')

        # Data rows
        approved_forms = [f for f in forms if f.estado.value == 'APROBADO']

        db = SessionLocal()
        try:
            crud = FormularioCRUD(db)

            for row, form in enumerate(approved_forms, 5):
                fresh_form = crud.get_formulario(form.id)
                if fresh_form:
                    ws.cell(row=row, column=1, value=form.id)
                    ws.cell(row=row, column=2, value=form.nombre_completo)
                    ws.cell(row=row, column=3, value=form.estado.value)
                    ws.cell(row=row, column=4, value=form.fecha_envio.strftime(
                        '%Y-%m-%d') if form.fecha_envio else '')
                    ws.cell(row=row, column=5, value=len(
                        fresh_form.cursos_capacitacion) if fresh_form.cursos_capacitacion else 0)
                    ws.cell(row=row, column=6, value=len(
                        fresh_form.publicaciones) if fresh_form.publicaciones else 0)
                    ws.cell(row=row, column=7, value=len(
                        fresh_form.eventos_academicos) if fresh_form.eventos_academicos else 0)
                    ws.cell(row=row, column=8, value=len(
                        fresh_form.certificaciones) if fresh_form.certificaciones else 0)
        finally:
            db.close()

        # Save to buffer
        buffer = BytesIO()
        wb.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    except ImportError:
        # Fallback to CSV-like content
        content = "ID,Docente,Estado,Fecha,Cursos,Publicaciones,Eventos,Certificaciones\n"
        approved_forms = [f for f in forms if f.estado.value == 'APROBADO']
        for form in approved_forms:
            content += f"{form.id},{form.nombre_completo},{form.estado.value},{form.fecha_envio.strftime('%Y-%m-%d') if form.fecha_envio else ''},0,0,0,0\n"
        return content.encode('utf-8')


def generate_powerpoint_report(forms, title, report_type, period_start, period_end):
    """Generate PowerPoint report using python-pptx"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from io import BytesIO

        # Create presentation
        prs = Presentation()

        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]

        title_placeholder.text = title
        subtitle_placeholder.text = f"Período: Año {period_start.year}" if period_start.year == period_end.year else f"Período: {period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}"

        # Summary slide
        slide_layout = prs.slide_layouts[1]  # Title and content
        slide = prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]

        title_placeholder.text = "Resumen de Actividades"

        # Calculate summary
        summary = calculate_activity_summary(forms)

        content_text = f"""Cursos de Capacitación: {summary['cursos']}
Publicaciones: {summary['publicaciones']}
Eventos Académicos: {summary['eventos']}
Diseños Curriculares: {summary['disenos']}
Experiencias de Movilidad: {summary['movilidades']}
Reconocimientos: {summary['reconocimientos']}
Certificaciones: {summary['certificaciones']}"""

        content_placeholder.text = content_text

        # Save to buffer
        buffer = BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()

    except ImportError:
        # Fallback to text content
        content = f"{title}\n\nPeríodo: Año {period_start.year}\n\n" if period_start.year == period_end.year else f"{title}\n\nPeríodo: {period_start.strftime('%B %Y')} - {period_end.strftime('%B %Y')}\n\n"
        content += generate_simple_report(forms, title,
                                          report_type, period_start, period_end)
        return content.encode('utf-8')
